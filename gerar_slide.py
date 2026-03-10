#!/usr/bin/env python3
"""
Gera apresentação PowerPoint - Parte 5: Sistema de Biblioteca Digital
Engenharia de Software - ADS
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Cores do tema
AZUL_ESCURO = RGBColor(0x1B, 0x2A, 0x4A)
AZUL_MEDIO = RGBColor(0x2C, 0x5F, 0x8A)
AZUL_CLARO = RGBColor(0x3A, 0x86, 0xC8)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO = RGBColor(0xF0, 0xF4, 0xF8)
CINZA_TEXTO = RGBColor(0x4A, 0x4A, 0x4A)
AMARELO = RGBColor(0xFF, 0xC1, 0x07)
VERDE = RGBColor(0x28, 0xA7, 0x45)
VERMELHO = RGBColor(0xDC, 0x35, 0x45)
LARANJA = RGBColor(0xFD, 0x7E, 0x14)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    """Preenche o fundo do slide com uma cor sólida."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=CINZA_TEXTO, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=CINZA_TEXTO, icon="▸", spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox


def add_card(slide, left, top, width, height, title, body, title_color=AZUL_MEDIO,
             bg_color=BRANCO, icon=""):
    """Adiciona um card estilizado."""
    shape = add_shape(slide, left, top, width, height, bg_color)
    shape.shadow.inherit = False

    # Título do card
    title_text = f"{icon}  {title}" if icon else title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15),
                width - Inches(0.5), Inches(0.5),
                title_text, font_size=16, color=title_color, bold=True)

    # Corpo do card
    add_textbox(slide, left + Inches(0.25), top + Inches(0.6),
                width - Inches(0.5), height - Inches(0.8),
                body, font_size=13, color=CINZA_TEXTO)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 - CAPA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, AZUL_ESCURO)

# Faixa decorativa superior
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), AMARELO)

# Faixa lateral esquerda
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, AZUL_CLARO)

# Ícone de livro estilizado (texto)
add_textbox(slide, Inches(5.2), Inches(1.2), Inches(3), Inches(1.2),
            "📚", font_size=72, color=BRANCO, alignment=PP_ALIGN.CENTER)

# Título principal
add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.5), Inches(1.2),
            "Sistema de Biblioteca Digital", font_size=44, color=BRANCO,
            bold=True, alignment=PP_ALIGN.CENTER)

# Subtítulo
add_textbox(slide, Inches(2.5), Inches(3.7), Inches(8.5), Inches(0.7),
            "Projeto de Desenvolvimento de Software", font_size=24,
            color=AZUL_CLARO, alignment=PP_ALIGN.CENTER)

# Linha divisória
add_shape(slide, Inches(5), Inches(4.5), Inches(3.333), Inches(0.04), AMARELO)

# Informações do curso
add_textbox(slide, Inches(2.5), Inches(4.9), Inches(8.5), Inches(0.5),
            "Engenharia de Software — Análise e Desenvolvimento de Sistemas",
            font_size=16, color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.5), Inches(5.4), Inches(8.5), Inches(0.5),
            "Professor: Danilson Ranniere Silva Ferro",
            font_size=14, color=RGBColor(0x88, 0x99, 0xBB), alignment=PP_ALIGN.CENTER)

# Rodapé
add_textbox(slide, Inches(2.5), Inches(6.4), Inches(8.5), Inches(0.5),
            "Caxias, Maranhão — Março de 2026",
            font_size=12, color=RGBColor(0x66, 0x77, 0x99), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 - AGENDA / SUMÁRIO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "📋  Agenda da Apresentação", font_size=32, color=BRANCO, bold=True)

agenda_items = [
    ("01", "Objetivo do Sistema", "O que é e para que serve a Biblioteca Digital"),
    ("02", "Usuários do Sistema", "Quem vai utilizar a plataforma"),
    ("03", "Funcionalidades Principais", "Recursos essenciais do sistema"),
    ("04", "Riscos do Projeto", "Ameaças ao sucesso do desenvolvimento"),
    ("05", "Etapas do Desenvolvimento", "Ciclo de vida do projeto"),
    ("06", "Problemas Potenciais", "Dificuldades que podem surgir"),
]

for i, (num, titulo, desc) in enumerate(agenda_items):
    y = Inches(1.5) + Inches(i * 0.9)
    # Número
    circle = add_shape(slide, Inches(1.2), y, Inches(0.6), Inches(0.6),
                       AZUL_MEDIO, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].font.color.rgb = BRANCO
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = "Calibri"

    # Título do item
    add_textbox(slide, Inches(2.1), y - Inches(0.05), Inches(5), Inches(0.4),
                titulo, font_size=20, color=AZUL_ESCURO, bold=True)
    # Descrição
    add_textbox(slide, Inches(2.1), y + Inches(0.3), Inches(8), Inches(0.35),
                desc, font_size=14, color=CINZA_TEXTO)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 - OBJETIVO DO SISTEMA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "🎯  Objetivo do Sistema", font_size=32, color=BRANCO, bold=True)

# Box principal
add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2), BRANCO)
add_textbox(slide, Inches(1.3), Inches(1.7), Inches(10.7), Inches(2.0),
            "Desenvolver uma plataforma digital completa para gerenciamento de acervo "
            "bibliográfico, permitindo que instituições de ensino e bibliotecas públicas "
            "ofereçam acesso remoto ao seu catálogo de livros, periódicos e materiais "
            "acadêmicos.\n\n"
            "O sistema visa modernizar o processo de empréstimo, devolução e consulta de "
            "obras, eliminando filas e burocracias, além de disponibilizar versões digitais "
            "(e-books e PDFs) para leitura online.",
            font_size=16, color=CINZA_TEXTO)

# Cards de benefícios
beneficios = [
    ("📖", "Acesso Remoto", "Consulta ao acervo\nde qualquer lugar"),
    ("⏱️", "Agilidade", "Empréstimos e devoluções\nsem filas"),
    ("🔍", "Busca Inteligente", "Pesquisa por título,\nautor ou assunto"),
    ("📊", "Relatórios", "Dados sobre uso\ndo acervo"),
]

for i, (icon, titulo, desc) in enumerate(beneficios):
    x = Inches(0.8) + Inches(i * 3.0)
    card = add_shape(slide, x, Inches(4.2), Inches(2.7), Inches(2.5), BRANCO)

    add_textbox(slide, x, Inches(4.3), Inches(2.7), Inches(0.7),
                icon, font_size=36, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(5.0), Inches(2.3), Inches(0.4),
                titulo, font_size=16, color=AZUL_MEDIO, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(5.4), Inches(2.3), Inches(0.8),
                desc, font_size=13, color=CINZA_TEXTO, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 - USUÁRIOS DO SISTEMA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "👥  Quem São os Usuários", font_size=32, color=BRANCO, bold=True)

usuarios = [
    ("🎓", "Alunos", "Estudantes de graduação e\npós-graduação que precisam\nconsultar e emprestar livros\npara seus estudos.", AZUL_CLARO),
    ("👨‍🏫", "Professores", "Docentes que necessitam\nreservar materiais para\ndisciplinas e indicar\nbibliografias aos alunos.", AZUL_MEDIO),
    ("📚", "Bibliotecários", "Profissionais responsáveis\npelo cadastro de obras,\ncontrole de acervo e\ngestão de empréstimos.", AZUL_ESCURO),
    ("🏛️", "Administradores", "Gestores do sistema com\nacesso a relatórios,\nconfigurações e controle\nde permissões.", RGBColor(0x5B, 0x21, 0xB6)),
]

for i, (icon, nome, desc, cor) in enumerate(usuarios):
    x = Inches(0.6) + Inches(i * 3.15)
    # Card
    card = add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(5.2), BRANCO)
    # Barra de cor no topo do card
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(0.08), cor)
    # Ícone
    add_textbox(slide, x, Inches(1.8), Inches(2.9), Inches(0.8),
                icon, font_size=48, alignment=PP_ALIGN.CENTER)
    # Nome
    add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(2.5), Inches(0.5),
                nome, font_size=22, color=cor, bold=True, alignment=PP_ALIGN.CENTER)
    # Descrição
    add_textbox(slide, x + Inches(0.25), Inches(3.3), Inches(2.4), Inches(3.0),
                desc, font_size=14, color=CINZA_TEXTO, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 - FUNCIONALIDADES (Parte 1)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "⚙️  Principais Funcionalidades", font_size=32, color=BRANCO, bold=True)

funcionalidades = [
    ("1", "Cadastro e Catalogação de Obras",
     "Registro completo de livros, periódicos, e-books e materiais acadêmicos com "
     "informações como título, autor, ISBN, editora, edição, ano e categoria. "
     "Suporte a código de barras e QR Code para identificação rápida."),
    ("2", "Sistema de Empréstimo e Devolução",
     "Controle automatizado de empréstimos com definição de prazos, renovações "
     "online e alertas de vencimento por e-mail/notificação. Multas automáticas "
     "para atrasos com possibilidade de pagamento online."),
    ("3", "Busca Avançada e Filtros",
     "Motor de busca com pesquisa por título, autor, palavra-chave, ISBN ou "
     "categoria. Filtros por disponibilidade, tipo de material, ano de publicação "
     "e área do conhecimento. Sugestões automáticas de obras relacionadas."),
    ("4", "Reserva Online de Livros",
     "Permite que o usuário reserve livros que estão emprestados, entrando em uma "
     "fila de espera com notificação automática quando o exemplar estiver "
     "disponível para retirada."),
]

for i, (num, titulo, desc) in enumerate(funcionalidades):
    y = Inches(1.4) + Inches(i * 1.45)
    # Número em círculo
    circle = add_shape(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.5),
                       AZUL_MEDIO, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].font.color.rgb = BRANCO
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = "Calibri"

    # Card
    add_shape(slide, Inches(1.5), y, Inches(11), Inches(1.3), BRANCO)
    add_textbox(slide, Inches(1.7), y + Inches(0.05), Inches(10.5), Inches(0.4),
                titulo, font_size=17, color=AZUL_ESCURO, bold=True)
    add_textbox(slide, Inches(1.7), y + Inches(0.45), Inches(10.5), Inches(0.8),
                desc, font_size=13, color=CINZA_TEXTO)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 - FUNCIONALIDADES (Parte 2)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "⚙️  Principais Funcionalidades (cont.)", font_size=32, color=BRANCO, bold=True)

funcionalidades2 = [
    ("5", "Leitura Digital (E-books / PDFs)",
     "Visualizador integrado para leitura de obras digitais diretamente na "
     "plataforma, com marcadores, anotações pessoais e modo noturno. "
     "Controle de DRM para proteger direitos autorais."),
    ("6", "Painel Administrativo e Relatórios",
     "Dashboard com estatísticas de uso: livros mais emprestados, horários de pico, "
     "usuários mais ativos e obras em atraso. Geração de relatórios em PDF para "
     "a gestão da instituição."),
    ("7", "Notificações e Alertas",
     "Sistema de avisos por e-mail e push sobre prazos de devolução, novas aquisições "
     "do acervo, disponibilidade de reservas e eventos da biblioteca."),
]

for i, (num, titulo, desc) in enumerate(funcionalidades2):
    y = Inches(1.4) + Inches(i * 1.5)
    circle = add_shape(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.5),
                       AZUL_MEDIO, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].font.color.rgb = BRANCO
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = "Calibri"

    add_shape(slide, Inches(1.5), y, Inches(11), Inches(1.3), BRANCO)
    add_textbox(slide, Inches(1.7), y + Inches(0.05), Inches(10.5), Inches(0.4),
                titulo, font_size=17, color=AZUL_ESCURO, bold=True)
    add_textbox(slide, Inches(1.7), y + Inches(0.45), Inches(10.5), Inches(0.8),
                desc, font_size=13, color=CINZA_TEXTO)

# Resumo visual
add_shape(slide, Inches(1.5), Inches(5.8), Inches(11), Inches(1.0), AZUL_ESCURO)
add_textbox(slide, Inches(1.7), Inches(5.9), Inches(10.5), Inches(0.8),
            "✅ Total: 7 funcionalidades principais que cobrem todo o ciclo de uso "
            "da biblioteca — do cadastro à leitura digital, passando por empréstimos, "
            "buscas, reservas e gestão administrativa.",
            font_size=15, color=BRANCO, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 - RISCOS DO PROJETO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "⚠️  Possíveis Riscos do Projeto", font_size=32, color=BRANCO, bold=True)

riscos = [
    ("🔴", "ALTO", "Segurança de Dados dos Usuários",
     "Vazamento de informações pessoais (CPF, endereço, histórico de leitura) "
     "pode gerar processos judiciais e violar a LGPD. É necessário implementar "
     "criptografia, autenticação robusta e backups regulares.",
     VERMELHO),
    ("🟠", "ALTO", "Resistência à Adoção por Parte dos Usuários",
     "Bibliotecários e usuários acostumados ao sistema manual podem resistir à "
     "migração digital. Treinamentos insuficientes podem levar ao abandono do "
     "sistema e retorno aos processos antigos.",
     LARANJA),
    ("🟡", "MÉDIO", "Indisponibilidade do Sistema (Downtime)",
     "Falhas de servidor, quedas de internet ou bugs podem tornar o acervo "
     "inacessível, prejudicando alunos em período de provas ou pesquisas. "
     "Plano de contingência e hospedagem redundante são essenciais.",
     AMARELO),
    ("🟡", "MÉDIO", "Escopo Mal Definido (Scope Creep)",
     "Solicitações constantes de novas funcionalidades durante o desenvolvimento "
     "podem atrasar entregas e estourar o orçamento. Gestão rigorosa de "
     "requisitos e controle de mudanças são fundamentais.",
     AMARELO),
    ("🟢", "BAIXO", "Incompatibilidade com Dispositivos",
     "O sistema pode não funcionar adequadamente em navegadores antigos ou "
     "dispositivos móveis. Testes de compatibilidade e design responsivo "
     "minimizam esse risco.",
     VERDE),
]

for i, (icon, nivel, titulo, desc, cor) in enumerate(riscos):
    y = Inches(1.3) + Inches(i * 1.18)
    # Barra de cor à esquerda
    add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(1.05), cor)
    # Card
    add_shape(slide, Inches(0.88), y, Inches(11.6), Inches(1.05), BRANCO)
    # Nível
    add_textbox(slide, Inches(1.1), y + Inches(0.05), Inches(1.0), Inches(0.3),
                f"{icon} {nivel}", font_size=11, color=cor, bold=True)
    # Título
    add_textbox(slide, Inches(1.1), y + Inches(0.3), Inches(3.0), Inches(0.3),
                titulo, font_size=15, color=AZUL_ESCURO, bold=True)
    # Descrição
    add_textbox(slide, Inches(4.0), y + Inches(0.05), Inches(8.3), Inches(0.95),
                desc, font_size=12, color=CINZA_TEXTO)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 - ETAPAS DO DESENVOLVIMENTO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "🔄  Etapas do Desenvolvimento", font_size=32, color=BRANCO, bold=True)

# Subtítulo - metodologia
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.5),
            "Metodologia: Modelo Incremental com práticas ágeis (Scrum)",
            font_size=16, color=AZUL_MEDIO, bold=True)

etapas = [
    ("1", "Levantamento\nde Requisitos",
     "Entrevistas com\nbibliotecários, alunos\ne professores",
     "2 semanas"),
    ("2", "Análise e\nModelagem",
     "Diagramas UML,\nprotótipos de tela\ne arquitetura",
     "3 semanas"),
    ("3", "Projeto\n(Design)",
     "Interface UI/UX,\nbanco de dados\ne APIs",
     "2 semanas"),
    ("4", "Implementação\n(Codificação)",
     "Desenvolvimento\nfrontend, backend\ne integração",
     "8 semanas"),
    ("5", "Testes",
     "Unitários, integração,\nusabilidade e\ndesempenho",
     "3 semanas"),
    ("6", "Implantação\ne Treinamento",
     "Deploy, migração\nde dados e\ncapacitação",
     "2 semanas"),
    ("7", "Manutenção\ne Evolução",
     "Correções, updates\ne novas features\ncontínuas",
     "Contínuo"),
]

for i, (num, nome, desc, tempo) in enumerate(etapas):
    x = Inches(0.35) + Inches(i * 1.85)
    y_base = Inches(1.9)

    # Seta conectora (exceto no último)
    if i < len(etapas) - 1:
        add_shape(slide, x + Inches(1.7), y_base + Inches(0.6),
                  Inches(0.3), Inches(0.04), AZUL_CLARO)
        # Ponta da seta
        add_shape(slide, x + Inches(1.85), y_base + Inches(0.5),
                  Inches(0.15), Inches(0.25), AZUL_CLARO,
                  MSO_SHAPE.CHEVRON)

    # Círculo do número
    circle = add_shape(slide, x + Inches(0.5), y_base, Inches(0.6), Inches(0.6),
                       AZUL_MEDIO, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(20)
    circle.text_frame.paragraphs[0].font.color.rgb = BRANCO
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.name = "Calibri"

    # Linha vertical
    add_shape(slide, x + Inches(0.77), y_base + Inches(0.6),
              Inches(0.06), Inches(0.4), AZUL_CLARO)

    # Card da etapa
    card = add_shape(slide, x + Inches(0.05), y_base + Inches(1.1),
                     Inches(1.55), Inches(3.2), BRANCO)

    # Nome da etapa
    add_textbox(slide, x + Inches(0.1), y_base + Inches(1.2),
                Inches(1.45), Inches(0.7),
                nome, font_size=13, color=AZUL_ESCURO, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Descrição
    add_textbox(slide, x + Inches(0.1), y_base + Inches(2.0),
                Inches(1.45), Inches(1.2),
                desc, font_size=11, color=CINZA_TEXTO,
                alignment=PP_ALIGN.CENTER)

    # Tempo
    tempo_shape = add_shape(slide, x + Inches(0.2), y_base + Inches(3.5),
                            Inches(1.2), Inches(0.35), AZUL_ESCURO,
                            MSO_SHAPE.ROUNDED_RECTANGLE)
    tempo_shape.text_frame.paragraphs[0].text = f"⏱ {tempo}"
    tempo_shape.text_frame.paragraphs[0].font.size = Pt(10)
    tempo_shape.text_frame.paragraphs[0].font.color.rgb = BRANCO
    tempo_shape.text_frame.paragraphs[0].font.bold = True
    tempo_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tempo_shape.text_frame.paragraphs[0].font.name = "Calibri"

# Tempo total
add_shape(slide, Inches(3.5), Inches(6.6), Inches(6.5), Inches(0.6), AZUL_ESCURO)
add_textbox(slide, Inches(3.5), Inches(6.65), Inches(6.5), Inches(0.5),
            "📅  Prazo estimado total: ~20 semanas (5 meses) + manutenção contínua",
            font_size=15, color=BRANCO, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 - PROBLEMAS QUE PODEM OCORRER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, CINZA_CLARO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), AZUL_ESCURO)
add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
            "🚧  Problemas que Podem Ocorrer", font_size=32, color=BRANCO, bold=True)

problemas = [
    ("Comunicação Falha na Equipe",
     "Mal-entendidos entre desenvolvedores, designers e o cliente podem gerar "
     "funcionalidades que não atendem à necessidade real. Reuniões periódicas "
     "(dailys e sprints) ajudam a mitigar esse problema.",
     "💬"),
    ("Mudanças Constantes de Requisitos",
     "O cliente pode solicitar alterações frequentes durante o desenvolvimento, "
     "gerando retrabalho e atrasos. Um contrato de escopo bem definido e "
     "processo formal de change requests são fundamentais.",
     "🔄"),
    ("Problemas de Integração com Sistemas Legados",
     "A biblioteca pode já possuir sistemas antigos (planilhas, bancos Access) "
     "cujos dados precisam ser migrados. Formatos incompatíveis e dados "
     "inconsistentes tornam a migração complexa e propensa a erros.",
     "🔗"),
    ("Subestimação de Prazos e Custos",
     "Equipes inexperientes podem estimar prazos otimistas demais, sem considerar "
     "imprevistos técnicos, curva de aprendizado e período de testes. Isso gera "
     "pressão excessiva e queda de qualidade.",
     "📉"),
    ("Falta de Testes Adequados",
     "Pular etapas de testes para cumprir prazos pode levar a bugs em produção, "
     "perdas de dados e frustração dos usuários. Automação de testes e QA "
     "dedicado são investimentos necessários.",
     "🐛"),
    ("Problemas de Desempenho sob Carga",
     "Em períodos de matrícula ou provas, muitos acessos simultâneos podem "
     "derrubar o sistema se não houver planejamento de escalabilidade, "
     "cache e otimização de consultas ao banco.",
     "📊"),
]

for i, (titulo, desc, icon) in enumerate(problemas):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.2)
    y = Inches(1.3) + Inches(row * 2.0)

    # Card
    add_shape(slide, x, y, Inches(5.9), Inches(1.8), BRANCO)
    # Ícone
    add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.6), Inches(0.6),
                icon, font_size=28)
    # Título
    add_textbox(slide, x + Inches(0.7), y + Inches(0.1), Inches(5.0), Inches(0.4),
                titulo, font_size=15, color=AZUL_ESCURO, bold=True)
    # Descrição
    add_textbox(slide, x + Inches(0.7), y + Inches(0.5), Inches(5.0), Inches(1.2),
                desc, font_size=12, color=CINZA_TEXTO)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 - CONCLUSÃO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AZUL_ESCURO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), AMARELO)

add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10.5), Inches(1.0),
            "✅  Conclusão", font_size=36, color=BRANCO, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(2.0), Inches(9.5), Inches(1.5),
            "O Sistema de Biblioteca Digital representa uma solução moderna e "
            "necessária para instituições de ensino que buscam otimizar a gestão "
            "do seu acervo e oferecer uma experiência digital de qualidade aos "
            "seus usuários.",
            font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), alignment=PP_ALIGN.CENTER)

# Destaques finais
destaques = [
    "7 funcionalidades essenciais mapeadas",
    "5 riscos identificados e classificados",
    "7 etapas de desenvolvimento planejadas",
    "6 problemas potenciais antecipados",
]

for i, item in enumerate(destaques):
    y = Inches(3.8) + Inches(i * 0.55)
    add_shape(slide, Inches(4.5), y, Inches(4.5), Inches(0.45), AZUL_MEDIO,
              MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, Inches(4.5), y + Inches(0.02), Inches(4.5), Inches(0.4),
                f"✔  {item}", font_size=15, color=BRANCO, bold=True,
                alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2), Inches(6.2), Inches(9.5), Inches(0.8),
            "\"Um bom software não nasce do código,\nnasce do planejamento.\"",
            font_size=16, color=AMARELO, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 - OBRIGADO
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AZUL_ESCURO)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), AMARELO)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.2),
            "Obrigado!", font_size=54, color=BRANCO, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.5), Inches(0.8),
            "📚 Sistema de Biblioteca Digital", font_size=24,
            color=AZUL_CLARO, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.333), Inches(0.04), AMARELO)

add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.5), Inches(0.5),
            "Dúvidas ou sugestões?", font_size=20,
            color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.5),
            "Engenharia de Software — ADS — Caxias/MA — 2026",
            font_size=14, color=RGBColor(0x66, 0x77, 0x99), alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SALVAR
# ═══════════════════════════════════════════════════════════════
output_path = "/home/paulo/Documentos/gambiarras/Parte5_Biblioteca_Digital.pptx"
prs.save(output_path)
print(f"✅ Apresentação salva em: {output_path}")
print(f"   Total de slides: {len(prs.slides)}")
